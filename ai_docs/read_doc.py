from fastapi import HTTPException, UploadFile
from pptx import Presentation
import PyPDF2, docx
import tempfile
import io, os
import win32com.client
import subprocess

async def process_file(file: UploadFile) -> str:
    """파일 확장자에 따라 파일을 읽고 텍스트를 추출합니다."""
    contents = file.file.read()
    if file.filename.endswith('.pdf'):
        return read_pdf(contents)
    elif file.filename.endswith('.docx'):
        return read_docx(contents)
    elif file.filename.endswith('.doc'):
        return read_doc(contents)
    elif file.filename.endswith('.pptx'):
        return read_pptx(contents)
    elif file.filename.endswith('.txt'):
        return read_txt(contents)
    else:
        raise HTTPException(status_code=400, detail="지원하지 않는 파일 형식입니다.")

def read_pdf(file_bytes):
    file_stream = io.BytesIO(file_bytes)
    reader = PyPDF2.PdfReader(file_stream)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text



def read_docx(file_bytes):
    file_stream = io.BytesIO(file_bytes)
    doc = docx.Document(file_stream)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def read_pptx(file_bytes):
    presentation = Presentation(io.BytesIO(file_bytes))
    
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return text

def read_doc(file_bytes):
    """DOC 파일을 읽고 텍스트를 추출합니다."""
    temp_file_path = os.path.join(tempfile.gettempdir(), "temp.doc")
    
    with open(temp_file_path, "wb") as temp_file:
        temp_file.write(file_bytes)

    try:
        word = win32com.client.Dispatch("Word.Application")
        doc = word.Documents.Open(temp_file_path)
        text = doc.Content.Text
        doc.Close()
        return text
    except Exception as e:
        raise RuntimeError(f"DOC 파일을 읽는 중 오류 발생: {e}")
    finally:
        os.remove(temp_file_path)

def read_txt(file_bytes):
    """TXT 파일을 읽고 텍스트를 추출합니다."""
    try:
        # 바이트를 문자열로 변환
        text = file_bytes.decode('utf-8')  # UTF-8 인코딩으로 디코딩
    except Exception as e:
        raise RuntimeError(f"TXT 파일 처리 중 오류 발생: {str(e)}")

    return text.strip()  # 공백 제거 후 반환

